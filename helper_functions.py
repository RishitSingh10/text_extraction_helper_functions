# -*- coding: utf-8 -*-
"""
Created on Tue Aug  6 13:00:03 2024

@author: rishit.somvanshi
"""
import streamlit as st
import translators as ts
import os
from pptx import Presentation
# from googletrans import Translator
import keras_ocr
import easyocr
# import pytesseract
import cv2
from pptx.enum.shapes import MSO_SHAPE_TYPE
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
import textwrap
import io
import numpy as np
import re
import warnings
from openai import OpenAI
from google.colab.patches import cv2_imshow
import tempfile
import pandas as pd
import matplotlib.pyplot as plt
from sklearn.cluster import DBSCAN
from sklearn import metrics
from sklearn.preprocessing import StandardScaler, MinMaxScaler
from sklearn.neighbors import NearestNeighbors
from kneed import KneeLocator

from bs4 import BeautifulSoup
import requests

warnings.filterwarnings("ignore")


def save_uploaded_file(uploaded_file):
    try:
        with open(os.path.join("/content", uploaded_file.name), "wb") as f:
            f.write(uploaded_file.getbuffer())
        return True
    except Exception as e:
        print(e)
        return False
    
def scrape_url(url, translator_service):
  response = requests.get(url)
  soup = BeautifulSoup(response.content, 'html.parser')
  tables = soup.find_all('table')
  translators = []

  if tables:
    first_table = tables[1]

    # Extract table rows
    rows = first_table.find_all('tr')
    headers = [service.get_text(strip=True) for service in rows[0].find_all('th')]
    translator_index = headers.index(translator_service)

    # List to hold data from the first two columns
    data = []

    for row in rows:
      columns = row.find_all('td')
      if len(columns) >= 2:
        # Get text from the first two columns
        lang = columns[0].get_text(strip=True)
        code = columns[1].get_text(strip=True)
        ifY = columns[translator_index].get_text(strip=True)

        if 'Y' in ifY:
          data.append((lang, code))

    return data

  else:
    return None
    
def translate_text(text, target_language):
  lang_code = [code for language, code in lang_code_data if language == target_language][0]
  translated_text = ''
  split_chars = ['\n', '.', '。', ',', '，']

  parts = re.split(r'(\n|\.|。|,|，)', text)  # Split using regex to preserve delimiters

  for part in parts:
    if part in split_chars:
      translated_text += part + ' '  # Append the delimiter directly
    else:
      translated_text += ts.translate_text(part, translator_service.lower(), to_language = lang_code)

  # print(f'{text} -------> {translated_text}')

  return translated_text


## TRANSLATE TXT FILES
def translate_txt(filename, target_language):

  tl = []

  with open(filename, 'r') as f:
    lines = f.readlines()
    for line in lines:
      translated_text = translate_text(line, target_language = target_language)
      tl.append(translated_text)

    f.close()

  with open(f'{target_language}_{filename}', 'w') as f:
    for line in tl:
      line = line.strip() + '\n'
      f.write(line)

    f.close()

  return f'{target_language}_{filename}'

## TRANSLATE IMAGES WITH TEXT

### Helper Functions for Text extraction and Clustering

# Function to get midpoint of bounding box
def midpoint(box):

  x_1, y_1 = box[2] # bottom right of first box
  x_2, y_2 = box[0] # bottom right of next box

  return [(x_1 + x_2)/2, (y_1 + y_2)/2]

# Function to get the avg height of bounding boxes
def get_avg_height(boxes):
  heights = []
  for box in boxes:
    x_1, y_1 = box[3] # top left
    x_2, y_2 = box[0] # bottom left

    heights.append(y_1 - y_2)

  return sum(heights)/len(heights)

# Function to use keras_ocr for text extraction
def keras_ocr_extraction(rgb_image):

  pipeline = keras_ocr.pipeline.Pipeline()
  prediction_groups = pipeline.recognize([rgb_image])

  return prediction_groups[0]

# Function to make a dataframe out of the predictions
def make_dataframe(text_box_pair):

  dataset = {}

  dataset['Word'] = [text for text, box in text_box_pair]
  dataset['x'] = [midpoint(box)[0] for _, box in text_box_pair]
  dataset['y'] = [midpoint(box)[1] for _, box in text_box_pair]
  boxes = [box for _, box in text_box_pair]

  df = pd.DataFrame(dataset)
  df = df.sort_values(by = 'x')

  df.reset_index(inplace = True, drop = True)

  return df, boxes, dataset

def find_optimal_eps(df_scaled, min_samples):

  neighbors = min_samples
  neigh = NearestNeighbors(n_neighbors=neighbors)
  nbrs = neigh.fit(df_scaled)
  distances, indices = nbrs.kneighbors(df_scaled)

  # Sort distances
  distances = np.sort(distances[:, min_samples-1])

  # Plot distances
  plt.plot(distances)
  plt.ylabel('Distance')
  plt.xlabel('Points sorted by distance')
  plt.title('K-distance Graph')
  plt.show()

  # Choose the eps based on the "elbow" point in the plot
  kneedle = KneeLocator(range(len(distances)), distances, curve='convex', direction='increasing')
  eps = kneedle.elbow_y
  print(f'Optimal eps: {eps}')

  return eps



# Function to cluster words so that each cluster can be treated as paragraph and translated together
def dbscan_clustering(df):

  scaler = MinMaxScaler()
  df_scaled = scaler.fit_transform(df[['x','y']])

  min_samples = 2
  eps = find_optimal_eps(df_scaled, min_samples)

  db = DBSCAN(eps=eps, min_samples=min_samples).fit(df_scaled)
  labels = db.labels_
  n_clusters_ = len(set(labels))
  core_samples_mask = np.zeros_like(db.labels_, dtype=bool)
  core_samples_mask[db.core_sample_indices_] = True

  return labels, n_clusters_

# Function to get a list of cluster of midpoints
def get_clusters_of_midpoints(labels, df, n_clusters_):

  midpoint_clusters = []

  for i in range(n_clusters_):
    midpoints = []
    for j, label in enumerate(labels):
      if label+1 == i:
        midpoints.append([df.iloc[j]['x'], df.iloc[j]['y']])

    midpoint_clusters.append(midpoints)

  return midpoint_clusters

# Function to group points by y-coordinate within a given tolerance
def group_as_lines(cluster, tolerance):
  cluster = np.array(cluster)
  sorted_points = cluster[np.argsort(cluster[:, 1])]
  lines = []

  current_group = [sorted_points[0]]
  for point in sorted_points[1:]:
      if abs(point[1] - current_group[-1][1]) <= tolerance:
          current_group.append(point)
      else:
          lines.append(current_group)
          current_group = [point]
  lines.append(current_group)

  return lines

# Function to sort each line in a cluster based on x axis
def get_sorted_clusters(midpoint_clusters, tolerance):

  cluster_as_lines = []

  for cluster in midpoint_clusters:
    if len(cluster) == 0:
      continue
    cluster_as_lines.append(group_as_lines(cluster, tolerance))

  sorted_clusters = []

  for cluster in cluster_as_lines:

    sorted_lines = []

    for line in cluster:
      line = np.array(line)
      line = line[np.argsort(line[:,0])]

      sorted_lines.append(line)

    sorted_clusters.append(sorted_lines)

  return sorted_clusters

def get_cluster_of_words(sorted_clusters, df):

  cluster_of_words = []

  for cluster in sorted_clusters:
    words = []
    for line in cluster:
      for point in line:
        words.append(df.iloc[np.where((df['x'] == point[0]) & (df['y'] == point[1]))[0][0]]['Word'])

    cluster_of_words.append(words)

  return cluster_of_words

def get_corrected_clusters(cluster_of_words):

  OPENAI_API_KEY = 'sk-proj-r9XeOGD9fwmwY0NjcDEyT3BlbkFJLPbp054JGa7Msrlusodv'
  client = OpenAI(api_key=OPENAI_API_KEY)

  corrected_clusters = []

  for list_ in cluster_of_words:
    prompt_2 = f"""You are an English Instructor. Here is a list of words: {list_}.
    If possible, try to make a sentence or a phrase or a paragraph without adding extra words.
    If it is not possible to make a sentence, just correct the words spelling.
    Either return the output only as a list or a sentence."""

    response = client.chat.completions.create(
      model="gpt-4o-mini",
      messages=[{"role": "user", "content":prompt_2}],
      temperature=1,
      max_tokens=256,
      top_p=1,
      frequency_penalty=0,
      presence_penalty=0
    )

    corrected_clusters.append(response.choices[0].message.content)

  return corrected_clusters

### Helper functions to put translated text back on the image

def get_box_label_df(dataset, labels, df):

  box_plus_label = []

  for (x, y) in zip(dataset['x'], dataset['y']):
    index = df[(df['x'] == x) & (df['y'] == y)].index[0]
    box_plus_label.append(labels[index])

  box_df = pd.DataFrame(box_plus_label, columns = ['label'])


  return box_df


def get_crop_colors(crop):

  with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
    temp_path = temp_file.name
    cv2.imwrite(temp_path, crop)

  with Image.open(temp_path) as i:

    color_map = i.getcolors() # colors is a list of tuples  with number of pixels containing  the colour and the RGB colour code.
    if not color_map:
      # Convert the image to palette mode
      img_palette = i.convert("P")

      # Get the colors in the image, allowing for up to 256 colors
      colors = img_palette.getcolors(256)

      # Get the palette of the image
      palette = img_palette.getpalette()

      # Map palette indices to RGB values
      color_map = []
      for count, index in colors:
          rgb = palette[index*3:index*3+3]  # Each color is represented by 3 values (R, G, B)
          color_map.append((count, tuple(rgb)))

    color_map = sorted(color_map)

    dominant_color = color_map[-1][1]
    weak_color = color_map[2][1]

  return dominant_color, weak_color

def calculate_font_scale(text, target_width, target_height, font=cv2.FONT_HERSHEY_SIMPLEX):
  # Start with a small font scale
  thickness= 2
  font_scale = 1
  max_font_scale = 10  # Arbitrary large number to avoid infinite loops

  while font_scale <= max_font_scale:
      # Calculate the size of the text with the current font scale
      (text_width, text_height), baseline = cv2.getTextSize(text, font, font_scale, thickness)

      # Check if the text fits within the target dimensions
      if text_width <= target_width and text_height <= target_height:
          # Scale is acceptable, but keep checking to maximize font size
          font_scale += 0.1
          thickness += 1
      else:
          # Text does not fit, so return the previous acceptable font scale
          return font_scale - 0.3, thickness-1

  return font_scale - 0.3, thickness-1

def get_min_max_coordinates(box_cluster_indexes, boxes):
  min_x = min_y = float('inf')
  max_x = max_y = float('-inf')

  for index in box_cluster_indexes:
    box = boxes[index]

    min_x = min(min_x, min(box[0][0], box[3][0]))
    min_y = min(min_y, min(box[0][1], box[1][1]))

    max_x = max(max_x, max(box[1][0], box[2][0]))
    max_y = max(max_y, max(box[2][1], box[3][1]))

  return (int(min_x), int(max_y)), (int(max_x), int(min_y))


def wrap_text(draw, text, font, max_width, text_x, text_y, weak_color):

  for line in textwrap.wrap(text, width=max_width):
    draw.text((text_x, text_y), line, font=font, fill=weak_color)
    text_y += font.getsize(line)[1]

def put_text_back(image, sorted_clusters, corrected_clusters, boxes, box_df, target_language):

  for i, cluster in enumerate(corrected_clusters):

    # To write other language text using PIL
    pil_image = Image.fromarray(image)
    draw = ImageDraw.Draw(pil_image)
    label = i-1

    if label == -1:

      for j, index in enumerate(box_df[box_df['label'] == label].index):

        box = boxes[index]
        # Get bounding box coordinates
        top_left = tuple(map(int, box[0]))
        bottom_right = tuple(map(int, box[2]))

        target_height = abs(bottom_right[1] - top_left[1])
        target_width = abs(bottom_right[0] - top_left[0])

        crop = image[top_left[1]:bottom_right[1], top_left[0]:bottom_right[0]]

        dominant_color, weak_color = get_crop_colors(cv2.cvtColor(crop, cv2.COLOR_RGB2BGR))

        try:
          lst = eval(cluster)
          text = translate_text(lst[j], target_language)
        except:
          text = translate_text(cluster, target_language)
          continue

        # To write other language text using PIL
        pil_image = Image.fromarray(image)
        draw = ImageDraw.Draw(pil_image)

        # Calculate text size
        font_scale, font_thickness = calculate_font_scale(text, target_width, target_height)

        # Use arial as it supports many characters
        font_path = "/content/Arial_Unicode_Font.ttf"
        font_size = int(font_scale*20)
        font = ImageFont.truetype(font_path, font_size, layout_engine=ImageFont.Layout.RAQM)

        draw.rectangle([top_left, bottom_right], fill = dominant_color)

        # Calculate the position to place the new text
        text_x = top_left[0]
        text_y = top_left[1]

        # Put the translated text on the image
        draw.text((text_x, text_y), text, font = font, fill = weak_color)

        image = np.array(pil_image)

    else:

      box_cluster_indexes = box_df[box_df['label'] == label].index

      (bottom_left, top_right) = get_min_max_coordinates(box_cluster_indexes, boxes)

      crop = image[top_right[1]:bottom_left[1], bottom_left[0]:top_right[0]]

      dominant_color, weak_color = get_crop_colors(cv2.cvtColor(crop, cv2.COLOR_RGB2BGR))
      draw.rectangle([bottom_left, top_right], fill = dominant_color)

      try:
        lst = eval(cluster)
        text = ' '.join(lst)
        text = translate_text(text, target_language)
      except:
        text = translate_text(cluster, target_language)


      target_width = int(abs(bottom_left[0] - top_right[0]))
      target_height = int(abs(bottom_left[1] - top_right[1]))

      # Calculate text size
      font_scale, font_thickness = calculate_font_scale(text, target_width, target_height)
      # (text_width, text_height), baseline = cv2.getTextSize(text, cv2.FONT_HERSHEY_SIMPLEX, font_scale, font_thickness)

      # Use arial as it supports many characters
      font_path = "/content/Arial_Unicode_Font.ttf"
      font_size = int(font_scale*20)
      font = ImageFont.truetype(font_path, font_size, layout_engine=ImageFont.Layout.RAQM)

      text_x, text_y = (bottom_left[0], top_right[1])
      target_width_wrap = int(target_width/font.getsize(text[0])[0]) # Number of characters in each line

      wrap_text(draw, text, font, target_width_wrap, text_x, text_y, weak_color)

      image = np.array(pil_image)

  cv2_imshow(image)
  cv2.waitKey(0)
  cv2.destroyAllWindows()

  return image

## TRANSLATE IMAGE FILES

def image_text_translation(image, target_language):

  text_box_pair = keras_ocr_extraction(image)

  df, boxes, dataset = make_dataframe(text_box_pair)
  tolerance = get_avg_height(boxes)/2

  labels, n_clusters_ = dbscan_clustering(df)

  midpoint_clusters = get_clusters_of_midpoints(labels, df, n_clusters_)
  sorted_clusters = get_sorted_clusters(midpoint_clusters, tolerance)
  cluster_of_words = get_cluster_of_words(sorted_clusters, df)
  corrected_clusters = get_corrected_clusters(cluster_of_words)

  box_df = get_box_label_df(dataset, labels, df)

  image = put_text_back(image, sorted_clusters, corrected_clusters, boxes, box_df, target_language)

  return image


## TRANSLATE PPT FILES

# Function to load ppt image shape as an image
def load_image(shape):
  image_shape = shape.image
  image_bytes = image_shape.blob

  # Load the image into a PIL Image object
  image_stream = io.BytesIO(image_bytes)
  image = Image.open(image_stream)
  image = image.convert('RGB')

  # Convert the PIL Image to an RGB numpy array (which Keras-OCR expects)
  rgb_image = np.array(image)

  return rgb_image

def translate_ppt(filename, target_language):
  prs = Presentation(filename)

  slides = prs.slides

  for slide in slides:

    if slide.has_notes_slide:
      notes_slide = slide.notes_slide
      text_frame = notes_slide.notes_text_frame
      text_frame.text = translate_text(text_frame.text, target_language)

    shapes = slide.shapes

    num_shapes = len(shapes)

    for i, shape in enumerate(shapes):

      # Check if the shape has text
      if hasattr(shape, "text"):

        font_name = ''
        font_size = int()
        bold = bool()

        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:

          alignment = paragraph.alignment

          for run in paragraph.runs:
            text = run.text
            translated_text = translate_text(text, target_language)

            font = run.font
            font_name = font.name
            font_size = font.size
            bold = font.bold

            run.text = translated_text
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = bold

          paragraph.alignment = alignment

      # Check if the shape is a table
      if shape.has_table:
        table = shape.table
        for row in table.rows:
          for cell in row.cells:
            text = cell.text
            translated_text = translate_text(text, target_language)
            cell.text = translated_text

      # Check if the shape is a picture
      if (shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
        image = load_image(shape)
        modified_image = image_text_translation(image, target_language)
        modified_image_pil = Image.fromarray(modified_image)  # Convert numpy array to PIL Image

        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
          temp_path = temp_file.name
          modified_image_pil.save(temp_path)
          shape._element.getparent().remove(shape._element)
          slide.shapes.add_picture(temp_path, shape.left, shape.top, shape.width, shape.height)


  prs.save(f'{filename[:-5]}_{target_language}.pptx')